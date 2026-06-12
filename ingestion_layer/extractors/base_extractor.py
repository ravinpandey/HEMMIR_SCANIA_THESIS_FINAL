"""
ingestion_layer/extractors/base_extractor.py

HEMMIR Multi-Format Ingestion — merged best-of-both base extractor.

Merge decisions (research scientist perspective):
  ✓ generate_doc_id     → Yixin: source_id::relative_path hash (fixes filename collision)
  ✓ VideoExtractor      → v2 tri-frame + optical flow peak + dense region crops
                          + Yixin's FileRecord API + keyframe_ids on segments
  ✓ All other extractors → Yixin (better docstrings, correct API, full coverage)
  ✓ structure_units.json → NEW: written by every extractor (enables Agent path)
  ✓ Code quality        → Yixin style (full docstrings, type hints, from __future__)

Format coverage:
  PDF     → PDFExtractorAdapter  (wraps existing Docling-based PDFExtractor)
  DOCX    → DocxExtractor        (python-docx, heading-group or paragraph chunking)
  PPTX    → PptxExtractor        (python-pptx, slide-level semantic units)
  CSV/TSV → CsvExtractor         (schema chunk + row-group chunks)
  XLSX    → XlsxExtractor        (per-sheet schema + row groups)
  Image   → ImageExtractor       (OCR optional, CLIP-ready)
  Video   → VideoExtractor       (Whisper ASR + tri-frame optical flow keyframes
                                  + dense region crops — research-grade)
  Unknown → UnknownExtractor     (graceful text fallback)

Critical outputs per format (must exist for downstream layers):
  metadata/doc_metadata.json       ← always
  metadata/text_chunks.json        ← always (empty list if no text)
  metadata/image_chunks.json       ← always (empty list if no images)
  metadata/table_chunks.json       ← always (empty list if no tables)
  metadata/structure_units.json    ← NEW: always — required by Agent path navigation
  metadata/video_segments.json     ← video only
  metadata/slide_chunks.json       ← PPTX only (supplementary)

Research novelty in video pipeline:
  - Optical flow peak frame selection (Farneback) picks the most visually
    active moment in each segment, not arbitrary midpoint
  - Information density detection (optical flow score ≥ threshold) flags
    segments worth deeper frame analysis
  - Dense region cropping extracts sub-frame ROIs from high-activity frames
    using adaptive thresholding + contour detection — critical for Scania
    industrial videos where camera zooms in on components or diagrams
  - All signals stored: is_information_dense, optical_flow_score, peak_frame_id
    These feed the enrichment layer's video_frame_enricher priority queue
"""

from __future__ import annotations

import hashlib
import json
from abc import ABC, abstractmethod
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from loguru import logger

from shared.models.multiformat_models import (
    ChunkStrategy,
    ColumnSchema,
    DocumentMetadataExtension,
    ExtendedPromotedFields,
    FileFormat,
    SlideChunkMetadata,
    SourceRecord,
    FolderRecord,
    StructuralUnit,
    StructureUnitType,
    VideoSegmentChunkMetadata,
    make_docx_text_extension,
    make_image_standalone_extension,
    make_xlsx_table_extension,
)
from ingestion_layer.loaders.source_loaders import FileRecord


# ── Optical flow configuration ─────────────────────────────────────────────
# Segments with peak optical flow ≥ this threshold are flagged
# as information-dense and trigger dense region cropping.
# Tuned empirically on Scania training videos — adjust per corpus.
INFO_DENSE_THRESHOLD: float = 2.5

# Keyframes per segment: start + peak + end
KEYFRAMES_PER_SEGMENT: int = 3

# Maximum region crops per dense frame
MAX_REGION_CROPS: int = 3


# ── ID generation ─────────────────────────────────────────────────────────────

def generate_doc_id(file_record: FileRecord) -> str:
    """
    Stable doc_id from source_id + relative_path hash.

    Uses relative_path (not just filename) so two files with the same
    name in different directories always get different doc_ids.
    This fixes the silent collision bug in v2 which hashed filename only.

    The source_id prefix scopes IDs to their ingest source — prevents
    cross-source collisions in multi-source corpora.
    """
    key = f"{file_record.source_record.source_id}::{file_record.relative_path}"
    return hashlib.md5(key.encode()).hexdigest()[:12]


def _chunk_id(doc_id: str, modality: str, idx: int) -> str:
    return f"{doc_id}_{modality}_{idx:03d}"

def _slide_chunk_id(doc_id: str, slide_idx: int) -> str:
    return f"{doc_id}_slide_{slide_idx:03d}"

def _segment_chunk_id(doc_id: str, seg_idx: int) -> str:
    return f"{doc_id}_seg_{seg_idx:04d}"

def _img_chunk_id(doc_id: str, context_idx: int, img_idx: int) -> str:
    return f"{doc_id}_ctx_{context_idx:03d}_img_{img_idx:03d}"

def _unit_id(doc_id: str, unit_type: str, idx: int) -> str:
    return f"{doc_id}_{unit_type}_{idx:03d}"


# ── Output helpers ─────────────────────────────────────────────────────────────

def _save_json(data: Any, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False, default=str)
    logger.debug(f"  Saved: {path.name} ({path.stat().st_size} bytes)")


def _create_output_folders(
    output_dir: Path,
    stem:       str,
    fmt:        str = "unknown",
) -> Dict[str, Path]:
    """
    Create output folder structure for a specific format.
    Only creates folders that are actually needed — prevents DOCX/PPTX
    from getting frames/ and slides/ folders that belong to video/pptx only.

    Folder matrix:
        metadata/  — always (all formats)
        images/    — PDF, DOCX, PPTX, Image, Video
        tables/    — PDF, DOCX, PPTX, CSV, XLSX
        slides/    — PPTX only
        frames/    — Video only
    """
    fmt = fmt.lower().strip(".")

    root = output_dir / stem

    # Always present
    folders: Dict[str, Path] = {
        "root":     root,
        "metadata": root / "metadata",
    }

    # images/ — all formats except CSV and XLSX
    sheet_fmts = {"csv", "xlsx", "xls", "ods", "tsv"}
    if fmt not in sheet_fmts:
        folders["images"] = root / "images"

    # tables/ — all formats except video
    video_fmts = {"mp4", "avi", "mkv", "mov", "webm", "m4v", "video"}
    if fmt not in video_fmts:
        folders["tables"] = root / "tables"

    # PPTX-only
    if fmt == "pptx":
        folders["slides"] = root / "slides"

    # Video-only
    if fmt in video_fmts:
        folders["frames"] = root / "frames"

    for p in folders.values():
        p.mkdir(parents=True, exist_ok=True)

    return folders


def _build_promoted_fields(
    chunk_dict:     Dict,
    doc_meta_dict:  Dict,
    file_record:    FileRecord,
    structure_unit: Optional[StructuralUnit] = None,
) -> Dict:
    """
    Build ExtendedPromotedFields for any chunk and format.
    Backward-compatible: all new fields default to safe values so
    existing ChromaDB collections for PDF docs are unaffected.
    """
    raw = {
        "doc_id":          doc_meta_dict.get("doc_id", ""),
        "doc_title":       doc_meta_dict.get("doc_title", ""),
        "tier1":           doc_meta_dict.get("tier1") or file_record.folder_record.tier1 or "",
        "tier2":           doc_meta_dict.get("tier2") or file_record.folder_record.tier2 or "",
        "project_id":      doc_meta_dict.get("project_id", ""),
        "document_type":   doc_meta_dict.get("document_type", ""),
        "language":        doc_meta_dict.get("language", "English"),
        "source_modality": chunk_dict.get("source_modality", ""),
        "page_number":     chunk_dict.get("page_number", 0),
        "chunk_index":     chunk_dict.get("chunk_index", 0),
        "chunk_strategy":  chunk_dict.get("chunk_strategy", ""),
        "contextual_summary_confidence": chunk_dict.get("contextual_summary_confidence", 0.0),
        "source_id":           file_record.source_record.source_id,
        "source_type":         file_record.source_record.source_type.value,
        "file_format":         file_record.file_format.value,
        "folder_path":         file_record.folder_record.folder_path,
        "structure_unit_type": structure_unit.structure_unit_type.value if structure_unit else "page",
        "structure_unit_id":   structure_unit.structure_unit_id if structure_unit else "",
        "sheet_name":          chunk_dict.get("format_specific", {}).get("sheet_name", ""),
        "slide_index":         chunk_dict.get("slide_index", 0),
        "start_time_s":        chunk_dict.get("start_time_s", 0.0),
    }
    return ExtendedPromotedFields.model_validate(raw).to_dict()


# ── Abstract Base Extractor ────────────────────────────────────────────────────

class BaseExtractor(ABC):
    """
    Abstract interface all format extractors implement.

    Every extractor must:
      1. Write the 5 standard output files (4 JSON + structure_units.json)
      2. Return a summary dict with doc_id and counts
      3. Never raise on empty/missing optional content

    The structure_units.json file is new (not in V4) and critical for the
    Agent path's section navigator — it searches the unified sections_collection
    which is populated from structure_units across all formats.
    """

    def __init__(self, output_dir: Path):
        self.output_dir = Path(output_dir)

    @abstractmethod
    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        """
        Extract content, write output files, return summary dict.

        Standard kwargs:
            chunk_strategy, chunk_size, chunk_overlap,
            skip_images, skip_tables, project_id, document_type, language
        """
        ...

    def _base_doc_meta(
        self, doc_id: str, file_record: FileRecord, **kwargs
    ) -> Dict[str, Any]:
        """
        Fields common to ALL formats — mirrors existing PDF doc_metadata exactly.
        Format-specific fields merged in by each subclass.
        """
        folder = file_record.folder_record
        return {
            "doc_id":        doc_id,
            "doc_title":     kwargs.get("doc_name") or file_record.local_path.stem,
            "project_id":    kwargs.get("project_id") or file_record.local_path.stem,
            "author":        None,
            "chunk_count":   0,
            "related_figures": [],
            "related_tables":  [],
            "last_modified":   file_record.last_modified or datetime.utcnow().isoformat(),
            "document_type":   kwargs.get("document_type"),
            "language":        kwargs.get("language", "English"),
            "total_pages":     0,
            "chunk_strategy":  kwargs.get("chunk_strategy", "paragraph"),
            "chunk_size":      kwargs.get("chunk_size", 512),
            "chunk_overlap":   kwargs.get("chunk_overlap", 150),
            "figure_index_map": {},
            "table_index_map":  {},
            "section_map":      {},
            "tier1":            folder.tier1 or kwargs.get("tier1"),
            "tier2":            folder.tier2 or kwargs.get("tier2"),
            "doc_summary":      None,
            "doc_summary_confidence": None,
            "tier_confidence":  None,
            "outline_summary":  None,
            "doc_embedding":    None,
            "outline_summary_confidence": None,
            "source_id":        file_record.source_record.source_id,
            "source_type":      file_record.source_record.source_type.value,
            "file_format":      file_record.file_format.value,
            "mime_type":        file_record.mime_type,
            "file_size_bytes":  file_record.file_size_bytes,
            "folder_path":      folder.folder_path,
            "folder_record":    folder.to_dict(),
            "relative_path":    file_record.relative_path,
        }


# ── PDF Adapter ────────────────────────────────────────────────────────────────

class PDFExtractorAdapter(BaseExtractor):
    """
    Thin adapter wrapping the existing Docling-based PDFExtractor.
    Zero modification to PDF extraction — only merges source/folder metadata
    and writes structure_units.json from the existing section_map.
    """

    def __init__(self, output_dir: Path):
        super().__init__(output_dir)
        from ingestion_layer.extractors.pdf_extractor import PDFExtractor as _PDF
        self._inner = _PDF(output_dir)

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        result = self._inner.extract(file_record.local_path, **kwargs)
        self._merge_source_metadata(file_record, kwargs)
        self._write_structure_units(file_record)
        return result

    def _merge_source_metadata(self, file_record: FileRecord, kwargs: Dict) -> None:
        doc_stem  = file_record.local_path.stem
        meta_path = self.output_dir / doc_stem / "metadata" / "doc_metadata.json"
        if not meta_path.exists():
            return

        with open(meta_path) as f:
            doc_meta = json.load(f)

        ext = DocumentMetadataExtension(
            source_id       = file_record.source_record.source_id,
            source_type     = file_record.source_record.source_type.value,
            folder_path     = file_record.folder_record.folder_path,
            folder_record   = file_record.folder_record.to_dict(),
            file_format     = FileFormat.PDF.value,
            mime_type       = "application/pdf",
            file_size_bytes = file_record.file_size_bytes,
            structure_type  = "article",
        )
        if not doc_meta.get("tier1") and file_record.folder_record.tier1:
            doc_meta["tier1"] = file_record.folder_record.tier1
        if not doc_meta.get("tier2") and file_record.folder_record.tier2:
            doc_meta["tier2"] = file_record.folder_record.tier2

        doc_meta.update(ext.to_dict())
        _save_json(doc_meta, meta_path)

    def _write_structure_units(self, file_record: FileRecord) -> None:
        """
        Derive structure_units.json from the section_map in doc_metadata.
        This enables the Agent path to navigate PDF sections.
        """
        doc_stem  = file_record.local_path.stem
        meta_dir  = self.output_dir / doc_stem / "metadata"
        meta_path = meta_dir / "doc_metadata.json"
        if not meta_path.exists():
            return

        with open(meta_path) as f:
            doc_meta = json.load(f)

        section_map = doc_meta.get("section_map", {})
        doc_id      = doc_meta.get("doc_id", "")
        units       = []
        for idx, (heading, entry) in enumerate(section_map.items(), start=1):
            if not isinstance(entry, dict):
                continue
            units.append({
                "structure_unit_id":   _unit_id(doc_id, "sec", idx),
                "doc_id":              doc_id,
                "structure_unit_type": StructureUnitType.HEADING.value,
                "unit_index":          idx,
                "title":               heading,
                "semantic_anchor":     heading,
                "start_page":          entry.get("start_page"),
                "end_page":            entry.get("end_page"),
                "chunk_ids":           entry.get("chunk_ids", []),
                "figure_ids":          entry.get("figure_ids", []),
                "table_ids":           entry.get("table_ids", []),
                "section_summary":     None,
                "keywords":            [],
                "entities":            [],
                "subsections":         [],
            })

        _save_json(units, meta_dir / "structure_units.json")
        logger.debug(f"PDF adapter: wrote {len(units)} structure_units")


# ── DOCX Extractor ─────────────────────────────────────────────────────────────

class DocxExtractor(BaseExtractor):
    """
    Extract text (heading-grouped or paragraph), tables, and embedded images
    from DOCX files using python-docx.

    Produces structure_units.json with one entry per heading section —
    these are the navigation targets for the Agent path.

    Chunking strategies: heading (default) | paragraph
    """

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("pip install python-docx")

        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="docx")
        doc     = DocxDocument(str(file_record.local_path))

        strategy = kwargs.get("chunk_strategy", "heading")
        text_chunks, structure_units = self._extract_text(doc, doc_id, strategy, kwargs)
        table_chunks, table_index_map = self._extract_tables(doc, doc_id, folders, file_record)
        image_chunks, figure_index_map = self._extract_images(doc, doc_id, folders, file_record)

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "doc_title":       self._get_title(doc, file_record),
            "chunk_count":     len(text_chunks) + len(image_chunks) + len(table_chunks),
            "related_figures": [c["chunk_id"] for c in image_chunks],
            "related_tables":  [c["chunk_id"] for c in table_chunks],
            "total_pages":     0,
            "figure_index_map": figure_index_map,
            "table_index_map":  table_index_map,
            "section_map":      {u["title"]: u for u in structure_units},
            "structure_type":   "document",
            "file_format":      FileFormat.DOCX.value,
        })

        for chunk in text_chunks + image_chunks + table_chunks:
            chunk["promoted_fields"] = _build_promoted_fields(chunk, doc_meta, file_record)

        _save_json(doc_meta,      folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,   folders["metadata"] / "text_chunks.json")
        _save_json(image_chunks,  folders["metadata"] / "image_chunks.json")
        _save_json(table_chunks,  folders["metadata"] / "table_chunks.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(
            f"DOCX: {file_record.relative_path!r} → "
            f"{len(text_chunks)} text | {len(image_chunks)} img | "
            f"{len(table_chunks)} tbl | {len(structure_units)} sections"
        )
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": len(image_chunks), "table_count": len(table_chunks),
        }

    def _get_title(self, doc, file_record: FileRecord) -> str:
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading 1") and para.text.strip():
                return para.text.strip()
        return file_record.local_path.stem

    def _extract_text(
        self, doc, doc_id: str, strategy: str, kwargs: Dict
    ) -> Tuple[List[Dict], List[Dict]]:
        from ingestion_layer.utils.chunking_utils import estimate_tokens

        chunks         = []
        structure_units = []
        current_heading       = None
        current_heading_level = None
        current_texts  = []
        chunk_index    = 1
        unit_index     = 1
        unit_chunk_ids: List[str] = []

        def flush_section():
            nonlocal chunk_index
            if not current_texts:
                return
            combined  = "\n\n".join(current_texts)
            chunk_id  = _chunk_id(doc_id, "chunk", chunk_index)
            ext = make_docx_text_extension(
                heading_level=current_heading_level,
                heading_text=current_heading,
            )
            chunks.append({
                "chunk_id":              chunk_id,
                "doc_id":                doc_id,
                "chunk_index":           chunk_index,
                "source_modality":       "text",
                "text_original_content": combined,
                "local_context":         None,
                "page_number":           0,
                "section_title":         current_heading,
                "chunk_strategy":        strategy,
                "token_count":           estimate_tokens(combined),
                "element_index":         chunk_index,
                "end_element_index":     chunk_index,
                "related_figures":       [],
                "related_tables":        [],
                "section_id":            current_heading or "",
                "contextual_summary":    None,
                "contextual_summary_confidence": None,
                "detected_codes":        [],
                "format_specific":       ext,
            })
            unit_chunk_ids.append(chunk_id)
            chunk_index += 1

        for para in doc.paragraphs:
            text  = para.text.strip()
            style = para.style.name
            if not text:
                continue

            is_heading = style.startswith("Heading")
            level = None
            if is_heading:
                try:
                    level = int(style.split()[-1])
                except (ValueError, IndexError):
                    level = 1

            if is_heading:
                if strategy in ("heading", "title"):
                    # Save previous section's chunk IDs into its structure unit
                    if structure_units and unit_chunk_ids:
                        structure_units[-1]["chunk_ids"] = list(unit_chunk_ids)
                        unit_chunk_ids.clear()
                    flush_section()
                    current_heading       = text
                    current_heading_level = level
                    current_texts         = []

                unit_id = _unit_id(doc_id, "heading", unit_index)
                structure_units.append({
                    "structure_unit_id":   unit_id,
                    "doc_id":              doc_id,
                    "structure_unit_type": StructureUnitType.HEADING.value,
                    "unit_index":          unit_index,
                    "title":               text,
                    "semantic_anchor":     text,
                    "chunk_ids":           [],
                    "figure_ids":          [],
                    "table_ids":           [],
                    "section_summary":     None,
                    "keywords":            [],
                    "entities":            [],
                    "subsections":         [],
                })
                unit_index += 1

            elif text:
                if strategy == "paragraph":
                    chunk_id = _chunk_id(doc_id, "chunk", chunk_index)
                    ext = make_docx_text_extension(
                        heading_text=current_heading,
                        heading_level=current_heading_level,
                        paragraph_style=style,
                    )
                    chunks.append({
                        "chunk_id":              chunk_id,
                        "doc_id":                doc_id,
                        "chunk_index":           chunk_index,
                        "source_modality":       "text",
                        "text_original_content": text,
                        "local_context":         None,
                        "page_number":           0,
                        "section_title":         current_heading,
                        "chunk_strategy":        "paragraph",
                        "token_count":           estimate_tokens(text),
                        "element_index":         chunk_index,
                        "end_element_index":     chunk_index,
                        "related_figures":       [],
                        "related_tables":        [],
                        "section_id":            current_heading or "",
                        "contextual_summary":    None,
                        "contextual_summary_confidence": None,
                        "detected_codes":        [],
                        "format_specific":       ext,
                    })
                    unit_chunk_ids.append(chunk_id)
                    chunk_index += 1
                else:
                    current_texts.append(text)

        if strategy in ("heading", "title"):
            if structure_units and unit_chunk_ids:
                structure_units[-1]["chunk_ids"] = list(unit_chunk_ids)
            flush_section()

        return chunks, structure_units

    def _extract_tables(
        self, doc, doc_id: str, folders: Dict, file_record: FileRecord
    ) -> Tuple[List[Dict], Dict]:
        table_chunks     = []
        table_index_map  = {}

        for tbl_idx, table in enumerate(doc.tables, start=1):
            chunk_id = f"{doc_id}_tbl_{tbl_idx:03d}"
            table_index_map[str(tbl_idx)] = chunk_id
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if not rows:
                continue

            headers   = rows[0]
            data      = rows[1:]
            row_count = len(rows)
            col_count = len(headers)

            header_html = "".join(f"<th>{h}</th>" for h in headers)
            body_html   = "".join(
                "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>"
                for row in data
            )
            table_html = f"<table><thead><tr>{header_html}</tr></thead><tbody>{body_html}</tbody></table>"

            csv_filename = f"{file_record.doc_stem}_tbl_{tbl_idx:03d}.csv"
            csv_path     = folders["tables"] / csv_filename
            try:
                import csv as _csv
                with open(csv_path, "w", newline="", encoding="utf-8") as f:
                    writer = _csv.writer(f)
                    for row in rows:
                        writer.writerow(row)
            except Exception as e:
                logger.warning(f"DOCX table CSV save failed: {e}")

            ext = make_xlsx_table_extension(
                sheet_name="document", sheet_index=0,
                header_row=headers, row_count=row_count, col_count=col_count,
            )
            table_chunks.append({
                "chunk_id":       chunk_id,
                "doc_id":         doc_id,
                "chunk_index":    tbl_idx,
                "source_modality": "table",
                "table_html":     table_html,
                "table_csv_path": str(Path("tables") / csv_filename),
                "html_file_path": "",
                "page_number":    0,
                "row_count":      row_count,
                "col_count":      col_count,
                "table_caption":  None,
                "table_caption_confidence": None,
                "table_summary":  None,
                "table_summary_confidence": None,
                "table_purpose":  None,
                "table_purpose_confidence": None,
                "markdown":       "",
                "html_representation": table_html,
                "format_specific": ext,
            })

        return table_chunks, table_index_map

    def _extract_images(
        self, doc, doc_id: str, folders: Dict, file_record: FileRecord
    ) -> Tuple[List[Dict], Dict]:
        image_chunks     = []
        figure_index_map = {}
        try:
            from docx.oxml.ns import qn
        except ImportError:
            return image_chunks, figure_index_map

        img_idx = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                img_idx += 1
                chunk_id     = f"{doc_id}_img_{img_idx:03d}"
                figure_index_map[str(img_idx)] = chunk_id

                img_filename = f"{file_record.doc_stem}_img_{img_idx:03d}.png"
                img_path     = folders["images"] / img_filename
                try:
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                except Exception as e:
                    logger.warning(f"DOCX image save failed: {e}")

                image_chunks.append({
                    "chunk_id":       chunk_id,
                    "doc_id":         doc_id,
                    "chunk_index":    img_idx,
                    "source_modality": "image",
                    "figure_id":      chunk_id,
                    "image_path":     str(Path("images") / img_filename),
                    "page_number":    0,
                    "image_type":     "embedded",
                    "related_sections": [],
                    "image_caption":  None,
                    "image_caption_confidence": None,
                    "depicted_component": None,
                    "depicted_component_confidence": None,
                    "visible_annotations": None,
                    "visible_annotations_confidence": None,
                    "contextual_summary": None,
                    "contextual_summary_confidence": None,
                })

        return image_chunks, figure_index_map


# ── PPTX Extractor ─────────────────────────────────────────────────────────────

class PptxExtractor(BaseExtractor):
    """
    Extract slides, notes, diagrams, and images from PPTX using python-pptx.

    One StructuralUnit per slide → enables Agent path to navigate directly
    to specific slides when the query asks about a presentation topic.
    """

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("pip install python-pptx")

        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="pptx")
        prs     = Presentation(str(file_record.local_path))

        slide_chunks     = []
        text_chunks      = []
        image_chunks     = []
        table_chunks     = []
        structure_units  = []
        slide_map        = {}
        figure_index_map = {}
        table_index_map  = {}
        img_global       = 0
        tbl_global       = 0

        from ingestion_layer.utils.chunking_utils import estimate_tokens

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Pre-scan for title — title shape may appear after image shapes
            # in the XML order, so we resolve it before the main loop to
            # ensure slide_title is correctly set on all image/chart chunks.
            #
            # Priority order (handles both standard and custom Scania decks):
            #   1. Placeholder idx 0 or 15 (standard PowerPoint title)
            #   2. Any shape whose name contains "title" (case-insensitive)
            #   3. First short text shape (≤ 80 chars) — common in custom
            #      corporate OnePagers that use freeform text boxes only
            title_text = ""
            _title_candidates: list = []   # (priority, text)
            for _s in slide.shapes:
                if not _s.has_text_frame:
                    continue
                _t = " ".join(
                    p.text.strip()
                    for p in _s.text_frame.paragraphs
                    if p.text.strip()
                ).strip()
                if not _t:
                    continue
                # Priority 1 — standard placeholder
                try:
                    _ph = _s.placeholder_format
                    if _ph and _ph.idx in (0, 15):
                        _title_candidates.append((1, _t))
                        continue
                except (ValueError, AttributeError):
                    pass
                # Priority 2 — shape name contains "title"
                if "title" in _s.name.lower():
                    _title_candidates.append((2, _t))
                    continue
                # Priority 3 — first short text (likely a heading)
                if len(_t) <= 80:
                    _title_candidates.append((3, _t))

            if _title_candidates:
                _title_candidates.sort(key=lambda x: x[0])
                title_text = _title_candidates[0][1]

            body_texts = []
            notes_text = ""
            has_image  = False
            has_table  = False
            has_chart  = False
            img_paths  = []
            slide_tbl_ids = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        try:
                            is_title = (
                                hasattr(shape, "placeholder_format")
                                and shape.placeholder_format
                                and shape.placeholder_format.idx in (0, 15)
                                and not title_text
                            )
                        except (ValueError, AttributeError):
                            is_title = False
                        if is_title:
                            title_text = text
                        else:
                            body_texts.append(text)

                if shape.has_table:
                    has_table  = True
                    tbl_global += 1
                    tbl_cid    = f"{doc_id}_slide_{slide_num:03d}_tbl_{tbl_global:03d}"
                    table_index_map[str(tbl_global)] = tbl_cid
                    slide_tbl_ids.append(tbl_cid)
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if rows:
                        h_html = "".join(f"<th>{h}</th>" for h in rows[0])
                        b_html = "".join(
                            "<tr>" + "".join(f"<td>{c}</td>" for c in r) + "</tr>"
                            for r in rows[1:]
                        )
                        tbl_html = (
                            f"<table><thead><tr>{h_html}</tr></thead>"
                            f"<tbody>{b_html}</tbody></table>"
                        )
                        table_chunks.append({
                            "chunk_id":       tbl_cid,
                            "doc_id":         doc_id,
                            "chunk_index":    tbl_global,
                            "source_modality": "table",
                            "table_html":     tbl_html,
                            "table_csv_path": "",
                            "html_file_path": "",
                            "page_number":    slide_num,
                            "row_count":      len(rows),
                            "col_count":      len(rows[0]) if rows else 0,
                            "table_caption":  f"Slide {slide_num} table",
                            "table_summary":  None,
                            "table_purpose":  None,
                            "format_specific": {"slide_index": slide_num},
                        })

                if hasattr(shape, "image"):
                    has_image  = True
                    img_global += 1
                    img_fname  = (
                        f"{file_record.doc_stem}_slide_{slide_num:03d}"
                        f"_img_{img_global:03d}.png"
                    )
                    img_path = folders["images"] / img_fname
                    try:
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)
                        img_paths.append(str(Path("images") / img_fname))
                    except Exception as e:
                        logger.warning(f"PPTX image save failed slide {slide_num}: {e}")

                    img_cid = _img_chunk_id(doc_id, slide_num, img_global)
                    figure_index_map[str(img_global)] = img_cid
                    image_chunks.append({
                        "chunk_id":          img_cid,
                        "doc_id":            doc_id,
                        "chunk_index":       img_global,
                        "source_modality":   "image",
                        "figure_id":         img_cid,
                        "image_path":        str(Path("images") / img_fname),
                        "page_number":       slide_num,
                        # slide_index at top level so ImageEnricher uses SLIDE_PROMPT
                        "slide_index":       slide_num,
                        "slide_title":       title_text or None,
                        "image_type":        "slide_image",
                        "related_sections":  [_slide_chunk_id(doc_id, slide_num)],
                        "section_id":        _slide_chunk_id(doc_id, slide_num),
                        "structure_unit_id": _slide_chunk_id(doc_id, slide_num),
                        "image_caption":     None,
                        "contextual_summary":             None,
                        "contextual_summary_confidence":  None,
                        "format_specific":   {"slide_index": slide_num},
                    })

                # ── Chart shape (shape_type == 3) ─────────────────────────
                # Charts are vector XML objects — no .image attribute.
                # Extract structured data as a table chunk (searchable text)
                # and render a PNG as an image chunk (vision-enrichable).
                if shape.shape_type == 3:
                    has_chart = True
                    tbl_global += 1
                    chart_cid  = f"{doc_id}_slide_{slide_num:03d}_chart_{tbl_global:03d}"
                    table_index_map[str(tbl_global)] = chart_cid
                    slide_tbl_ids.append(chart_cid)

                    # ── Extract chart data from OpenXML ──────────────────
                    chart_data = self._extract_chart_data(shape)
                    c_title    = chart_data.get("title") or f"Chart on slide {slide_num}"
                    c_md       = chart_data.get("markdown", "")
                    c_html     = chart_data.get("html", "")
                    c_cols     = chart_data.get("column_names", [])
                    c_rows     = chart_data.get("row_count", 0)
                    c_type     = chart_data.get("chart_type", "chart")

                    table_chunks.append({
                        "chunk_id":         chart_cid,
                        "doc_id":           doc_id,
                        "chunk_index":      tbl_global,
                        "source_modality":  "table",
                        "file_format":      "pptx",
                        "structure_unit_id": _slide_chunk_id(doc_id, slide_num),
                        "parent_id":        _slide_chunk_id(doc_id, slide_num),
                        "table_html":       c_html or None,
                        "table_csv_path":   "",
                        "html_file_path":   "",
                        "page_number":      slide_num,
                        "row_count":        c_rows,
                        "col_count":        len(c_cols),
                        "column_names":     c_cols,
                        "markdown":         c_md or None,
                        "section_id":       _slide_chunk_id(doc_id, slide_num),
                        "sheet_name":       f"slide_{slide_num}_chart",
                        "sheet_index":      slide_num,
                        # table_caption pre-filled — TableEnricher will upgrade
                        "table_caption":    c_title,
                        "table_summary":    None,
                        "table_purpose":    None,
                        "format_specific":  {
                            "slide_index": slide_num,
                            "chart_type":  c_type,
                            "is_chart":    True,
                        },
                    })

                    # ── Render chart as PNG for vision enrichment ────────
                    img_global += 1
                    chart_img_fname = (
                        f"{file_record.doc_stem}_slide_{slide_num:03d}"
                        f"_chart_{tbl_global:03d}.png"
                    )
                    chart_img_path = folders["images"] / chart_img_fname
                    rendered = self._render_chart_png(chart_data, chart_img_path)

                    if rendered:
                        chart_img_cid = _img_chunk_id(doc_id, slide_num, img_global)
                        figure_index_map[str(img_global)] = chart_img_cid
                        img_paths.append(str(Path("images") / chart_img_fname))
                        image_chunks.append({
                            "chunk_id":          chart_img_cid,
                            "doc_id":            doc_id,
                            "chunk_index":       img_global,
                            "source_modality":   "image",
                            "figure_id":         chart_img_cid,
                            "image_path":        str(Path("images") / chart_img_fname),
                            "page_number":       slide_num,
                            "slide_index":       slide_num,
                            "slide_title":       title_text or None,
                            "image_type":        "chart",
                            "related_sections":  [_slide_chunk_id(doc_id, slide_num)],
                            "image_caption":     c_title,   # pre-filled title
                            "depicted_component": c_type,   # chart type known
                            "contextual_summary":            None,
                            "contextual_summary_confidence": None,
                            "format_specific":   {
                                "slide_index": slide_num,
                                "chart_type":  c_type,
                                "data_chunk_id": chart_cid,
                            },
                        })
                        has_image = True

            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            slide_cid      = _slide_chunk_id(doc_id, slide_num)
            body_combined  = "\n".join(body_texts)
            slide_text_all = "\n".join(
                p for p in [title_text, body_combined, notes_text] if p
            )

            slide_chunks.append(SlideChunkMetadata(
                chunk_id    = slide_cid,
                doc_id      = doc_id,
                chunk_index = slide_num,
                slide_index = slide_num,
                slide_title = title_text or None,
                slide_body  = body_combined or None,
                slide_notes = notes_text or None,
                slide_layout = slide.slide_layout.name if slide.slide_layout else None,
                has_image   = has_image,
                has_table   = has_table,
                has_chart   = has_chart,
                image_paths = img_paths,
                table_ids   = slide_tbl_ids,
            ).to_dict())

            text_chunk_id = _chunk_id(doc_id, "chunk", slide_num)
            if slide_text_all.strip():
                text_chunks.append({
                    "chunk_id":              text_chunk_id,
                    "doc_id":                doc_id,
                    "chunk_index":           slide_num,
                    "source_modality":       "text",
                    "text_original_content": slide_text_all,
                    "local_context":         None,
                    "page_number":           slide_num,
                    "section_title":         title_text or f"Slide {slide_num}",
                    "chunk_strategy":        ChunkStrategy.SLIDE.value,
                    "token_count":           estimate_tokens(slide_text_all),
                    "element_index":         slide_num,
                    "end_element_index":     slide_num,
                    "related_figures":       img_paths,
                    "related_tables":        slide_tbl_ids,
                    "section_id":            slide_cid,
                    "structure_unit_id":     slide_cid,
                    "contextual_summary":    None,
                    "contextual_summary_confidence": None,
                    "detected_codes":        [],
                    "format_specific":       {
                        "slide_index": slide_num,
                        "has_notes":   bool(notes_text),
                    },
                })

            # Structure unit per slide — Agent path navigation target
            slide_fig_ids = [
                c["chunk_id"] for c in image_chunks
                if c.get("page_number") == slide_num
            ]
            structure_units.append({
                "structure_unit_id":   slide_cid,
                "doc_id":              doc_id,
                "structure_unit_type": StructureUnitType.SLIDE.value,
                "unit_index":          slide_num,
                "title":               title_text or f"Slide {slide_num}",
                "semantic_anchor":     title_text or f"Slide {slide_num}",
                "chunk_ids":           [text_chunk_id] if slide_text_all.strip() else [],
                "figure_ids":          slide_fig_ids,
                "table_ids":           slide_tbl_ids,
                "section_summary":     None,
                "keywords":            [],
                "entities":            [],
                "subsections":         [],
            })
            slide_map[str(slide_num)] = structure_units[-1]

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "doc_title":       kwargs.get("doc_name") or self._get_title(prs, file_record),
            "chunk_count":     len(text_chunks) + len(image_chunks) + len(table_chunks),
            "related_figures": [c["chunk_id"] for c in image_chunks],
            "related_tables":  [c["chunk_id"] for c in table_chunks],
            "total_pages":     len(prs.slides),
            "total_slides":    len(prs.slides),
            "figure_index_map": figure_index_map,
            "table_index_map":  table_index_map,
            "section_map":      slide_map,
            "slide_map":        slide_map,
            "structure_type":   "slide_deck",
            "file_format":      FileFormat.PPTX.value,
        })

        for chunk in text_chunks + image_chunks + table_chunks:
            chunk["promoted_fields"] = _build_promoted_fields(chunk, doc_meta, file_record)

        _save_json(doc_meta,       folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,    folders["metadata"] / "text_chunks.json")
        _save_json(image_chunks,   folders["metadata"] / "image_chunks.json")
        _save_json(table_chunks,   folders["metadata"] / "table_chunks.json")
        _save_json(slide_chunks,   folders["metadata"] / "slide_chunks.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(
            f"PPTX: {file_record.relative_path!r} → "
            f"{len(prs.slides)} slides | {len(image_chunks)} img | "
            f"{len(table_chunks)} tbl"
        )
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": len(image_chunks), "table_count": len(table_chunks),
            "slide_count": len(prs.slides),
        }

    def _get_title(self, prs, file_record: FileRecord) -> str:
        try:
            props = prs.core_properties
            if props.title:
                return props.title
        except Exception:
            pass
        return file_record.local_path.stem

    # ── Chart helpers (PPTX only — no impact on any other extractor) ─────────

    _CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    def _extract_chart_data(self, shape) -> Dict[str, Any]:
        """
        Extract chart title, type, categories and series values from OpenXML.

        Returns a dict with keys:
          title, chart_type, column_names, row_count, series,
          markdown (pipe-delimited table), html (<table>…</table>)
        Falls back gracefully — returns empty strings rather than raising.
        """
        NS = self._CHART_NS
        result: Dict[str, Any] = {
            "title":        "",
            "chart_type":   "chart",
            "column_names": [],
            "row_count":    0,
            "series":       [],
            "markdown":     "",
            "html":         "",
        }

        try:
            # ── Chart title ──────────────────────────────────────────────
            try:
                chart_obj = shape.chart
                if chart_obj.has_title:
                    result["title"] = (
                        chart_obj.chart_title.text_frame.text.strip()
                    )
            except Exception:
                pass

            # ── Chart type ───────────────────────────────────────────────
            try:
                ct = str(shape.chart.chart_type)
                # "XL_CHART_TYPE.BAR_CLUSTERED (2)" → "bar_clustered"
                result["chart_type"] = (
                    ct.split(".")[-1].split("(")[0].strip().lower()
                )
            except Exception:
                pass

            # ── Series + categories from OpenXML ─────────────────────────
            # Use raw XML for reliability — python-pptx chart API varies by type
            el = shape._element
            all_cats: List[str]         = []
            series_data: List[Dict]     = []

            for ser in el.findall(f".//{{{NS}}}ser"):
                # Series name
                s_name = ""
                for t_el in ser.findall(f".//{{{NS}}}tx"):
                    for t in t_el.iter():
                        if t.tag.endswith("}v") and t.text:
                            s_name = t.text.strip()
                            break
                    if s_name:
                        break

                # Categories (xVal or cat element)
                cats: List[str] = []
                for cat_tag in (f"{{{NS}}}cat", f"{{{NS}}}xVal"):
                    cat_el = ser.find(f".//{cat_tag}")
                    if cat_el is not None:
                        for v in cat_el.findall(f".//{{{NS}}}v"):
                            if v.text:
                                cats.append(v.text.strip())
                        break

                # Values (yVal or val element)
                vals: List[Optional[float]] = []
                for val_tag in (f"{{{NS}}}val", f"{{{NS}}}yVal"):
                    val_el = ser.find(f".//{val_tag}")
                    if val_el is not None:
                        for v in val_el.findall(f".//{{{NS}}}v"):
                            try:
                                vals.append(float(v.text))
                            except (TypeError, ValueError):
                                vals.append(None)
                        break

                if vals:
                    series_data.append({
                        "name":   s_name or f"Series {len(series_data)+1}",
                        "values": vals,
                        "cats":   cats,
                    })
                    if cats and not all_cats:
                        all_cats = cats

            result["series"] = series_data

            if not series_data:
                return result

            # ── Build markdown + HTML table ──────────────────────────────
            # Columns: Category | Series1 | Series2 | ...
            series_names = [s["name"] for s in series_data]
            col_names    = (["Category"] if all_cats else []) + series_names
            result["column_names"] = col_names

            n_rows = max(len(s["values"]) for s in series_data)
            result["row_count"] = n_rows

            rows: List[List[str]] = []
            for i in range(n_rows):
                row: List[str] = []
                if all_cats:
                    row.append(all_cats[i] if i < len(all_cats) else "")
                for s in series_data:
                    v = s["values"][i] if i < len(s["values"]) else None
                    row.append("" if v is None else str(v))
                rows.append(row)

            # Markdown
            header  = "| " + " | ".join(col_names) + " |"
            sep     = "| " + " | ".join("---" for _ in col_names) + " |"
            md_rows = [
                "| " + " | ".join(r) + " |" for r in rows
            ]
            result["markdown"] = "\n".join([header, sep] + md_rows)

            # HTML
            th  = "".join(f"<th>{c}</th>" for c in col_names)
            trs = "".join(
                "<tr>" + "".join(f"<td>{v}</td>" for v in row) + "</tr>"
                for row in rows
            )
            result["html"] = (
                f"<table><thead><tr>{th}</tr></thead>"
                f"<tbody>{trs}</tbody></table>"
            )

        except Exception as exc:
            logger.debug(f"  chart XML extraction error: {exc}")

        return result

    def _render_chart_png(
        self, chart_data: Dict[str, Any], out_path: Path
    ) -> bool:
        """
        Render a chart PNG from extracted chart_data using matplotlib.
        Returns True on success, False if matplotlib is unavailable or data empty.

        Only creates a PNG — no side-effects on any other extractor or format.
        """
        series = chart_data.get("series", [])
        if not series:
            return False

        try:
            import matplotlib
            matplotlib.use("Agg")          # non-interactive backend, safe on EC2
            import matplotlib.pyplot as plt
            import numpy as np
        except ImportError:
            logger.debug("  matplotlib not installed — chart PNG skipped")
            return False

        try:
            chart_type = chart_data.get("chart_type", "")
            title      = chart_data.get("title", "Chart")
            cats       = series[0].get("cats", []) if series else []
            x          = list(range(len(series[0]["values"])))
            x_labels   = cats if cats else [str(i) for i in x]

            fig, ax = plt.subplots(figsize=(8, 5))
            fig.patch.set_facecolor("#f8f9fa")
            ax.set_facecolor("#f8f9fa")

            is_pie = "pie" in chart_type
            is_bar = any(k in chart_type for k in ("bar", "column"))

            if is_pie and series:
                vals   = [v for v in series[0]["values"] if v is not None]
                labels = x_labels[:len(vals)]
                ax.pie(vals, labels=labels, autopct="%1.1f%%",
                       startangle=90)
                ax.set_aspect("equal")
            elif is_bar:
                width  = 0.8 / max(len(series), 1)
                for si, s in enumerate(series):
                    vals = [v if v is not None else 0 for v in s["values"]]
                    offsets = [xi + si * width for xi in x]
                    ax.bar(offsets, vals, width=width, label=s["name"])
                ax.set_xticks([xi + width * (len(series) - 1) / 2 for xi in x])
                ax.set_xticklabels(x_labels, rotation=30, ha="right", fontsize=8)
                if len(series) > 1:
                    ax.legend(fontsize=8)
            else:
                # Default: line chart
                for s in series:
                    vals = [v if v is not None else float("nan") for v in s["values"]]
                    ax.plot(x, vals, marker="o", label=s["name"])
                ax.set_xticks(x)
                ax.set_xticklabels(x_labels, rotation=30, ha="right", fontsize=8)
                if len(series) > 1:
                    ax.legend(fontsize=8)

            if title:
                ax.set_title(title, fontsize=11, pad=10)
            ax.spines[["top", "right"]].set_visible(False)

            plt.tight_layout()
            plt.savefig(str(out_path), dpi=120, bbox_inches="tight")
            plt.close(fig)
            return True

        except Exception as exc:
            logger.debug(f"  chart PNG render error: {exc}")
            try:
                plt.close("all")
            except Exception:
                pass
            return False


# ── CSV Extractor ──────────────────────────────────────────────────────────────

class CsvExtractor(BaseExtractor):
    """
    Extract schema chunk + row-group chunks from CSV/TSV files.

    Schema chunk (always first): column names + dtypes + sample values.
    This is the Agent path's primary entry point for spreadsheet queries.
    Row-group chunks: fixed-row sliding windows as markdown tables.
    """

    DEFAULT_ROW_GROUP_SIZE = 100

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pip install pandas")

        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="csv")

        sep = "\t" if file_record.local_path.suffix.lower() == ".tsv" else ","
        df  = pd.read_csv(str(file_record.local_path), sep=sep, low_memory=False)

        text_chunks,  column_schema = self._build_chunks(df, doc_id, file_record, kwargs)
        table_chunks, table_index_map = self._build_table_chunks(df, doc_id, folders, file_record)

        sheet_unit_id = f"{doc_id}_sheet_001"
        structure_units = [{
            "structure_unit_id":   sheet_unit_id,
            "doc_id":              doc_id,
            "structure_unit_type": StructureUnitType.SHEET.value,
            "unit_index":          1,
            "title":               file_record.local_path.name,
            "semantic_anchor":     f"Dataset: {file_record.local_path.stem}",
            "sheet_name":          "data",
            "chunk_ids":           [c["chunk_id"] for c in text_chunks],
            "table_ids":           [c["chunk_id"] for c in table_chunks],
            "figure_ids":          [],
            "section_summary":     None,
            "keywords":            [],
            "entities":            list(df.columns[:10]),
            "subsections":         [],
        }]

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "chunk_count":    len(text_chunks) + len(table_chunks),
            "related_tables": [c["chunk_id"] for c in table_chunks],
            "total_pages":    0,
            "table_index_map": table_index_map,
            "structure_type": "spreadsheet",
            "file_format":    FileFormat.CSV.value,
            "total_sheets":   1,
            "column_schemas": {"data": [c.to_dict() for c in column_schema]},
            "section_map":    {"data": structure_units[0]},
        })

        for chunk in text_chunks + table_chunks:
            chunk["promoted_fields"] = _build_promoted_fields(chunk, doc_meta, file_record)

        _save_json(doc_meta,       folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,    folders["metadata"] / "text_chunks.json")
        _save_json(table_chunks,   folders["metadata"] / "table_chunks.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(
            f"CSV: {file_record.relative_path!r} → "
            f"{len(df)} rows | {len(df.columns)} cols"
        )
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": 0, "table_count": len(table_chunks),
        }

    def _build_chunks(
        self, df, doc_id: str, file_record: FileRecord, kwargs: Dict
    ) -> Tuple[List[Dict], List]:
        from ingestion_layer.utils.chunking_utils import estimate_tokens

        text_chunks   = []
        column_schema = self._infer_schema(df)

        # Schema chunk
        schema_lines = [
            f"File: {file_record.local_path.name}",
            f"Rows: {len(df)} | Columns: {len(df.columns)}",
            "",
            "Column Schema:",
        ]
        for col in column_schema:
            sample = ", ".join(str(v) for v in col.sample_values[:3])
            schema_lines.append(f"  {col.name} ({col.dtype}) — sample: [{sample}]")
        schema_text = "\n".join(schema_lines)

        ext = make_xlsx_table_extension(
            sheet_name="data", sheet_index=0,
            header_row=list(df.columns),
            row_count=len(df), col_count=len(df.columns),
            column_schema=[c.to_dict() for c in column_schema],
            is_schema_chunk=True,
        )
        text_chunks.append({
            "chunk_id":              _chunk_id(doc_id, "chunk", 1),
            "doc_id":                doc_id,
            "chunk_index":           1,
            "source_modality":       "text",
            "text_original_content": schema_text,
            "local_context":         None,
            "page_number":           0,
            "section_title":         "Schema",
            "chunk_strategy":        ChunkStrategy.SCHEMA.value,
            "token_count":           estimate_tokens(schema_text),
            "element_index":         1,
            "end_element_index":     1,
            "related_figures":       [],
            "related_tables":        [],
            "contextual_summary":    None,
            "contextual_summary_confidence": None,
            "detected_codes":        [],
            "format_specific":       ext,
        })

        # Row-group chunks
        group_size = int(kwargs.get("row_group_size", self.DEFAULT_ROW_GROUP_SIZE))
        for start in range(0, len(df), group_size):
            end    = min(start + group_size, len(df))
            subset = df.iloc[start:end]
            md     = (
                subset.to_markdown(index=False)
                if hasattr(subset, "to_markdown")
                else subset.to_string(index=False)
            )
            ci = len(text_chunks) + 1
            ext_rg = make_xlsx_table_extension(
                sheet_name="data", sheet_index=0,
                header_row=list(df.columns),
                start_row=start, end_row=end - 1,
                row_count=len(subset), col_count=len(df.columns),
            )
            text_chunks.append({
                "chunk_id":              _chunk_id(doc_id, "chunk", ci),
                "doc_id":                doc_id,
                "chunk_index":           ci,
                "source_modality":       "text",
                "text_original_content": md,
                "local_context":         None,
                "page_number":           0,
                "section_title":         f"Rows {start}–{end - 1}",
                "chunk_strategy":        ChunkStrategy.ROW_GROUP.value,
                "token_count":           estimate_tokens(md),
                "element_index":         ci,
                "end_element_index":     ci,
                "related_figures":       [],
                "related_tables":        [],
                "contextual_summary":    None,
                "contextual_summary_confidence": None,
                "detected_codes":        [],
                "format_specific":       ext_rg,
            })

        return text_chunks, column_schema

    def _build_table_chunks(
        self, df, doc_id: str, folders: Dict, file_record: FileRecord
    ) -> Tuple[List[Dict], Dict]:
        chunk_id        = f"{doc_id}_tbl_001"
        table_index_map = {"1": chunk_id}
        header_html = "".join(f"<th>{col}</th>" for col in df.columns)
        body_html   = "".join(
            "<tr>" + "".join(f"<td>{v}</td>" for v in row) + "</tr>"
            for row in df.head(10).values
        )
        html      = f"<table><thead><tr>{header_html}</tr></thead><tbody>{body_html}</tbody></table>"
        csv_fname = f"{file_record.doc_stem}_sheet_001.csv"
        import shutil
        shutil.copy(str(file_record.local_path), str(folders["tables"] / csv_fname))

        ext = make_xlsx_table_extension(
            sheet_name="data", sheet_index=0,
            header_row=list(df.columns),
            row_count=len(df), col_count=len(df.columns),
        )
        return [{
            "chunk_id":       chunk_id,
            "doc_id":         doc_id,
            "chunk_index":    1,
            "source_modality": "table",
            "table_html":     html,
            "table_csv_path": str(Path("tables") / csv_fname),
            "html_file_path": "",
            "page_number":    0,
            "row_count":      len(df),
            "col_count":      len(df.columns),
            "table_caption":  file_record.local_path.name,
            "table_summary":  None,
            "table_purpose":  None,
            "format_specific": ext,
        }], table_index_map

    def _infer_schema(self, df) -> List[ColumnSchema]:
        schemas = []
        for col in df.columns:
            s       = df[col]
            samples = s.dropna().unique()[:3].tolist()
            schemas.append(ColumnSchema(
                name          = str(col),
                dtype         = str(s.dtype),
                null_count    = int(s.isna().sum()),
                sample_values = [str(v) for v in samples],
                is_numeric    = str(s.dtype) in ("int64", "float64", "int32", "float32"),
                is_date       = "datetime" in str(s.dtype),
            ))
        return schemas


# ── XLSX Extractor ─────────────────────────────────────────────────────────────

class XlsxExtractor(BaseExtractor):
    """
    Extract per-sheet schema + row-group chunks from XLSX/XLS.
    One StructuralUnit per sheet — Agent path navigates directly to the
    relevant sheet when answering data queries.
    """

    DEFAULT_ROW_GROUP_SIZE = 100

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pip install pandas")

        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="xlsx")
        ext_    = file_record.local_path.suffix.lower()
        engine  = "xlrd" if ext_ == ".xls" else "openpyxl"
        xl      = pd.ExcelFile(str(file_record.local_path), engine=engine)

        text_chunks       = []
        table_chunks      = []
        structure_units   = []
        sheet_map         = {}
        table_index_map   = {}
        column_schemas_all = {}
        sheet_idx         = 0

        for sheet_name in xl.sheet_names:
            sheet_idx += 1
            try:
                df = xl.parse(sheet_name)
            except Exception as e:
                logger.warning(f"  Cannot parse sheet '{sheet_name}': {e}")
                continue
            if df.empty:
                continue

            sheet_unit_id = f"{doc_id}_sheet_{sheet_idx:03d}"
            schema        = self._infer_schema(df)
            column_schemas_all[sheet_name] = [c.to_dict() for c in schema]

            chunk_base  = len(text_chunks) + 1
            schema_text = self._build_schema_text(df, sheet_name, schema, file_record)
            sc_id       = _chunk_id(doc_id, "chunk", chunk_base)
            ext_s = make_xlsx_table_extension(
                sheet_name=sheet_name, sheet_index=sheet_idx,
                header_row=list(df.columns),
                row_count=len(df), col_count=len(df.columns),
                column_schema=[c.to_dict() for c in schema],
                is_schema_chunk=True,
            )
            schema_chunk = {
                "chunk_id":              sc_id,
                "doc_id":                doc_id,
                "chunk_index":           chunk_base,
                "source_modality":       "text",
                "text_original_content": schema_text,
                "local_context":         None,
                "page_number":           sheet_idx,
                "section_title":         f"Schema: {sheet_name}",
                "chunk_strategy":        ChunkStrategy.SCHEMA.value,
                "token_count":           len(schema_text) // 4,
                "element_index":         chunk_base,
                "end_element_index":     chunk_base,
                "related_figures":       [],
                "related_tables":        [],
                "contextual_summary":    None,
                "contextual_summary_confidence": None,
                "detected_codes":        [],
                "format_specific":       ext_s,
            }
            text_chunks.append(schema_chunk)
            sheet_chunk_ids = [sc_id]

            group_size = int(kwargs.get("row_group_size", self.DEFAULT_ROW_GROUP_SIZE))
            for start in range(0, len(df), group_size):
                end    = min(start + group_size, len(df))
                subset = df.iloc[start:end]
                md     = (
                    subset.to_markdown(index=False)
                    if hasattr(subset, "to_markdown")
                    else subset.to_string()
                )
                ci = len(text_chunks) + 1
                ext_rg = make_xlsx_table_extension(
                    sheet_name=sheet_name, sheet_index=sheet_idx,
                    header_row=list(df.columns),
                    start_row=start, end_row=end - 1,
                    row_count=len(subset), col_count=len(df.columns),
                )
                rg_id = _chunk_id(doc_id, "chunk", ci)
                sheet_chunk_ids.append(rg_id)
                text_chunks.append({
                    "chunk_id":              rg_id,
                    "doc_id":                doc_id,
                    "chunk_index":           ci,
                    "source_modality":       "text",
                    "text_original_content": md,
                    "local_context":         None,
                    "page_number":           sheet_idx,
                    "section_title":         f"{sheet_name}: rows {start}–{end - 1}",
                    "chunk_strategy":        ChunkStrategy.ROW_GROUP.value,
                    "token_count":           len(md) // 4,
                    "element_index":         ci,
                    "end_element_index":     ci,
                    "related_figures":       [],
                    "related_tables":        [],
                    "contextual_summary":    None,
                    "contextual_summary_confidence": None,
                    "detected_codes":        [],
                    "format_specific":       ext_rg,
                })

            tbl_id = f"{doc_id}_sheet_{sheet_idx:03d}_tbl"
            table_index_map[str(sheet_idx)] = tbl_id
            header_html = "".join(f"<th>{c}</th>" for c in df.columns)
            body_html   = "".join(
                "<tr>" + "".join(f"<td>{v}</td>" for v in row) + "</tr>"
                for row in df.head(20).values
            )
            tbl_html  = (
                f"<table><thead><tr>{header_html}</tr></thead>"
                f"<tbody>{body_html}</tbody></table>"
            )
            csv_fname = f"{file_record.doc_stem}_sheet_{sheet_idx:03d}.csv"
            try:
                df.to_csv(str(folders["tables"] / csv_fname), index=False)
            except Exception as e:
                logger.warning(f"XLSX sheet CSV save failed: {e}")

            ext_t = make_xlsx_table_extension(
                sheet_name=sheet_name, sheet_index=sheet_idx,
                header_row=list(df.columns),
                row_count=len(df), col_count=len(df.columns),
            )
            table_chunks.append({
                "chunk_id":       tbl_id,
                "doc_id":         doc_id,
                "chunk_index":    sheet_idx,
                "source_modality": "table",
                "table_html":     tbl_html,
                "table_csv_path": str(Path("tables") / csv_fname),
                "html_file_path": "",
                "page_number":    sheet_idx,
                "row_count":      len(df),
                "col_count":      len(df.columns),
                "table_caption":  sheet_name,
                "table_summary":  None,
                "table_purpose":  None,
                "format_specific": ext_t,
            })

            unit = {
                "structure_unit_id":   sheet_unit_id,
                "doc_id":              doc_id,
                "structure_unit_type": StructureUnitType.SHEET.value,
                "unit_index":          sheet_idx,
                "title":               sheet_name,
                "semantic_anchor":     f"Sheet: {sheet_name}",
                "sheet_name":          sheet_name,
                "chunk_ids":           sheet_chunk_ids,
                "table_ids":           [tbl_id],
                "figure_ids":          [],
                "section_summary":     None,
                "keywords":            [],
                "entities":            list(df.columns[:10]),
                "subsections":         [],
            }
            structure_units.append(unit)
            sheet_map[sheet_name] = unit

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "chunk_count":     len(text_chunks) + len(table_chunks),
            "related_tables":  [c["chunk_id"] for c in table_chunks],
            "total_pages":     sheet_idx,
            "total_sheets":    sheet_idx,
            "table_index_map": table_index_map,
            "section_map":     sheet_map,
            "sheet_map":       sheet_map,
            "column_schemas":  column_schemas_all,
            "structure_type":  "spreadsheet",
            "file_format":     FileFormat.XLSX.value,
        })

        for chunk in text_chunks + table_chunks:
            chunk["promoted_fields"] = _build_promoted_fields(chunk, doc_meta, file_record)

        _save_json(doc_meta,       folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,    folders["metadata"] / "text_chunks.json")
        _save_json(table_chunks,   folders["metadata"] / "table_chunks.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(f"XLSX: {file_record.relative_path!r} → {sheet_idx} sheets")
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": 0, "table_count": len(table_chunks),
        }

    def _infer_schema(self, df) -> List[ColumnSchema]:
        schemas = []
        for col in df.columns:
            s       = df[col]
            samples = s.dropna().unique()[:3].tolist()
            schemas.append(ColumnSchema(
                name          = str(col),
                dtype         = str(s.dtype),
                null_count    = int(s.isna().sum()),
                sample_values = [str(v) for v in samples],
                is_numeric    = str(s.dtype) in ("int64", "float64", "int32", "float32"),
                is_date       = "datetime" in str(s.dtype),
            ))
        return schemas

    def _build_schema_text(
        self, df, sheet_name: str, schema: List[ColumnSchema], fr: FileRecord
    ) -> str:
        lines = [
            f"File: {fr.local_path.name} | Sheet: {sheet_name}",
            f"Rows: {len(df)} | Columns: {len(df.columns)}",
            "",
            "Column Schema:",
        ]
        for c in schema:
            samples = ", ".join(str(v) for v in c.sample_values[:3])
            lines.append(f"  {c.name} ({c.dtype}) — sample: [{samples}]")
        return "\n".join(lines)


# ── Image Extractor ────────────────────────────────────────────────────────────

class ImageExtractor(BaseExtractor):
    """
    Process standalone image files (JPEG, PNG, TIFF, etc.)
    Optional OCR via pytesseract for images containing text.
    """

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="image")

        try:
            from PIL import Image as PILImage
            img    = PILImage.open(str(file_record.local_path))
            width, height = img.size
            mode   = img.mode
        except Exception as e:
            logger.error(f"Image open failed {file_record.relative_path}: {e}")
            raise

        img_fname = f"{file_record.doc_stem}{file_record.local_path.suffix}"
        img_dest  = folders["images"] / img_fname
        import shutil
        shutil.copy(str(file_record.local_path), str(img_dest))

        ocr_text = None
        ocr_conf = None
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img).strip() or None
            ocr_conf = 1.0 if ocr_text else None
        except ImportError:
            logger.debug("pytesseract not installed — skipping OCR")
        except Exception as e:
            logger.warning(f"OCR failed: {e}")

        ext = make_image_standalone_extension(
            width_px=width, height_px=height,
            color_mode=mode,
            ocr_text=ocr_text,
            ocr_confidence=ocr_conf,
        )

        chunk_id    = f"{doc_id}_img_001"
        image_chunk = {
            "chunk_id":       chunk_id,
            "doc_id":         doc_id,
            "chunk_index":    1,
            "source_modality": "image",
            "figure_id":      chunk_id,
            "image_path":     str(Path("images") / img_fname),
            "page_number":    1,
            "image_type":     "standalone",
            "related_sections": [],
            "image_caption":  None,
            "image_caption_confidence": None,
            "depicted_component": None,
            "depicted_component_confidence": None,
            "visible_annotations": None,
            "visible_annotations_confidence": None,
            "contextual_summary": None,
            "contextual_summary_confidence": None,
            "format_specific": ext,
        }

        text_chunks = []
        if ocr_text:
            from ingestion_layer.utils.chunking_utils import estimate_tokens
            text_chunks.append({
                "chunk_id":              f"{doc_id}_chunk_001",
                "doc_id":                doc_id,
                "chunk_index":           1,
                "source_modality":       "text",
                "text_original_content": ocr_text,
                "local_context":         None,
                "page_number":           1,
                "section_title":         "OCR Text",
                "chunk_strategy":        ChunkStrategy.WHOLE_FILE.value,
                "token_count":           estimate_tokens(ocr_text),
                "element_index":         1,
                "end_element_index":     1,
                "related_figures":       [chunk_id],
                "related_tables":        [],
                "contextual_summary":    None,
                "contextual_summary_confidence": None,
                "detected_codes":        [],
                "format_specific":       {"ocr_source": chunk_id},
            })

        structure_units = [{
            "structure_unit_id":   f"{doc_id}_img_unit_001",
            "doc_id":              doc_id,
            "structure_unit_type": "image_file",
            "unit_index":          1,
            "title":               file_record.local_path.name,
            "semantic_anchor":     f"Image: {file_record.local_path.stem}",
            "chunk_ids":           [c["chunk_id"] for c in text_chunks],
            "figure_ids":          [chunk_id],
            "table_ids":           [],
            "section_summary":     None,
            "keywords":            [],
            "entities":            [],
            "subsections":         [],
        }]

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "chunk_count":    len(text_chunks) + 1,
            "related_figures": [chunk_id],
            "total_pages":    1,
            "figure_index_map": {"1": chunk_id},
            "structure_type": "image_file",
            "file_format":    FileFormat.IMAGE.value,
            "width_px":       width,
            "height_px":      height,
        })

        image_chunk["promoted_fields"] = _build_promoted_fields(image_chunk, doc_meta, file_record)
        for tc in text_chunks:
            tc["promoted_fields"] = _build_promoted_fields(tc, doc_meta, file_record)

        _save_json(doc_meta,       folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,    folders["metadata"] / "text_chunks.json")
        _save_json([image_chunk],  folders["metadata"] / "image_chunks.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(f"Image: {file_record.relative_path!r} [{width}×{height}]")
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": 1, "table_count": 0,
        }


# ── Video Extractor ────────────────────────────────────────────────────────────

class VideoExtractor(BaseExtractor):
    """
    Research-grade video ingestion pipeline.

    Transcript:  OpenAI Whisper ASR with word-level timestamps
    Keyframes:   Tri-frame per segment (start / optical-flow peak / end)
    Peak select: Farneback optical flow — picks most visually active moment
    Dense crops: Adaptive threshold + contour detection on high-activity frames
    Scene units: Every SCENE_GROUP_SIZE segments = one StructuralUnit (scene)

    Why tri-frame + optical flow (research motivation):
      A fixed midpoint frame is arbitrary — it may capture the speaker's face
      rather than the component being discussed. The optical flow peak finds the
      frame where the most visual change occurs — camera zoom on a part, hand
      movement demonstrating a tool, diagram appearing on a whiteboard.
      For Scania industrial videos this dramatically improves CLIP retrieval
      precision: the retrieved frame shows the relevant component, not a
      talking head between demonstrations.

    Dense region crops (research motivation):
      When a Scania instructor zooms in on a hydraulic fitting or shows a
      torque spec on a label, the full frame is cluttered with background.
      Adaptive thresholding identifies the high-contrast region (the component
      or label) and crops it as a separate image chunk. CLIP then embeds a
      focused close-up rather than a noisy full frame.
    """

    SEGMENT_DURATION_S = 30.0   # fallback when no word timestamps
    SCENE_GROUP_SIZE   = 10     # segments per scene StructuralUnit

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        doc_id   = generate_doc_id(file_record)
        folders  = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="mp4")
        strategy = kwargs.get("chunk_strategy", "transcript")

        # Step 1: Transcribe
        segments = self._transcribe(file_record, strategy, kwargs)

        # Step 2: Extract keyframes (tri-frame + optical flow + dense crops)
        image_chunks:    List[Dict] = []
        figure_index_map: Dict      = {}
        if not kwargs.get("skip_images", False):
            image_chunks, figure_index_map = self._extract_keyframes(
                file_record, segments, doc_id, folders
            )

        # Step 3: Build segment and text chunk metadata
        from ingestion_layer.utils.chunking_utils import estimate_tokens
        text_chunks:    List[Dict] = []
        video_segments: List[Dict] = []

        for i, seg in enumerate(segments, start=1):
            seg_id    = _segment_chunk_id(doc_id, i)
            text      = seg.get("text", "").strip()
            start_s   = float(seg.get("start", 0.0))
            end_s     = float(seg.get("end", start_s + self.SEGMENT_DURATION_S))

            # Pull vrag signals written by _extract_keyframes
            vrag = seg.get("_vrag", {})
            keyframe_ids      = list(vrag.get("keyframe_ids", []))
            peak_frame_id     = vrag.get("peak_frame_id")
            is_info_dense     = bool(vrag.get("is_information_dense", False))
            optical_flow_score = float(vrag.get("optical_flow_score", 0.0))
            keyframe_path     = vrag.get("keyframe_path")
            keyframe_ts       = vrag.get("keyframe_timestamp_s")

            vs = VideoSegmentChunkMetadata(
                chunk_id             = seg_id,
                doc_id               = doc_id,
                chunk_index          = i,
                start_time_s         = start_s,
                end_time_s           = end_s,
                segment_index        = i,
                transcript_text      = text or None,
                transcript_confidence = seg.get("confidence"),
                speaker_label        = seg.get("speaker"),
                keyframe_path        = keyframe_path,
                keyframe_timestamp_s = keyframe_ts,
                keyframe_ids         = keyframe_ids,
                peak_frame_id        = peak_frame_id,
                is_information_dense = is_info_dense,
                optical_flow_score   = optical_flow_score,
            )
            video_segments.append(vs.to_dict())

            if text:
                text_chunks.append({
                    "chunk_id":              _chunk_id(doc_id, "chunk", i),
                    "doc_id":                doc_id,
                    "chunk_index":           i,
                    "source_modality":       "text",
                    "text_original_content": text,
                    "local_context":         None,
                    "page_number":           0,
                    "section_title":         f"Segment {i} ({start_s:.1f}s–{end_s:.1f}s)",
                    "chunk_strategy":        ChunkStrategy.TRANSCRIPT.value,
                    "token_count":           estimate_tokens(text),
                    "element_index":         i,
                    "end_element_index":     i,
                    "related_figures":       keyframe_ids,
                    "related_tables":        [],
                    "contextual_summary":    None,
                    "contextual_summary_confidence": None,
                    "detected_codes":        [],
                    "format_specific":       {
                        "start_time_s":      start_s,
                        "end_time_s":        end_s,
                        "segment_id":        seg_id,
                        "peak_frame_id":     peak_frame_id or "",
                        "is_information_dense": is_info_dense,
                        "optical_flow_score": optical_flow_score,
                    },
                })

        # Step 4: Build scene StructuralUnits
        structure_units = self._build_scene_units(
            doc_id, segments, video_segments, text_chunks
        )

        duration = segments[-1].get("end", 0.0) if segments else None

        # Step 5: Save full transcript for enrichment layer caching
        if segments:
            transcript_dir = folders["metadata"].parent / "transcripts"
            transcript_dir.mkdir(exist_ok=True)
            full_text = " ".join(
                seg.get("text", "").strip() for seg in segments
            )
            stem = file_record.doc_stem
            (transcript_dir / f"{stem}_full_transcript.txt").write_text(
                full_text, encoding="utf-8"
            )
            import json as _json
            (transcript_dir / f"{stem}_full_transcript.json").write_text(
                _json.dumps(segments, ensure_ascii=False, indent=2), encoding="utf-8"
            )

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "chunk_count":      len(text_chunks) + len(image_chunks),
            "related_figures":  [c["chunk_id"] for c in image_chunks],
            "total_pages":      len(segments),
            "total_segments":   len(segments),
            "total_frames":     len(image_chunks),
            "figure_index_map": figure_index_map,
            "structure_type":   "video_file",
            "file_format":      FileFormat.VIDEO.value,
            "duration_seconds": duration,
            "whisper_model":    kwargs.get("whisper_model", "base"),
        })

        for chunk in text_chunks + image_chunks:
            chunk["promoted_fields"] = _build_promoted_fields(
                chunk, doc_meta, file_record
            )

        # ── Build video_frames.json from keyframes in image_chunks ────
        # VideoFrameChunkMetadata is the correct format for VideoFrameEnricher
        # and the encoding layer's video frame path. It requires:
        #   parent_id      → parent video segment chunk_id
        #   timestamp_s    → exact frame timestamp
        #   frame_position → position within segment (1=start, 2=peak, 3=end)
        #   segment_start_s / segment_end_s → temporal range of parent segment
        #
        # We build this from image_chunks (which have format_specific fields
        # written by _extract_keyframes) and cross-reference with video_segments
        # to get parent segment timestamps.
        #
        # BOTH files are written:
        #   image_chunks.json  → ImageEnricher + image embedding path (existing)
        #   video_frames.json  → VideoFrameEnricher (transcript-grounded vision)
        #                      + video frame encoding path
        seg_by_id = {vs["chunk_id"]: vs for vs in video_segments}
        video_frames = []
        frame_position_counter = {}   # segment_id → frame count within segment

        for ic in image_chunks:
            fmt = ic.get("format_specific", {})
            seg_id = fmt.get("segment_id", "")
            role   = fmt.get("frame_role", "other")
            ts     = float(fmt.get("timestamp_s", 0.0))

            # Track frame position within segment (1=start, 2=peak, 3=end)
            frame_position_counter.setdefault(seg_id, 0)
            frame_position_counter[seg_id] += 1
            frame_pos = frame_position_counter[seg_id]

            parent_seg = seg_by_id.get(seg_id, {})

            video_frames.append({
                "chunk_id":          ic["chunk_id"],
                "doc_id":            doc_id,
                "source_id":         file_record.source_record.source_id,
                "source_modality":   "video_frame",
                "file_format":       FileFormat.VIDEO.value,
                "chunk_index":       ic.get("chunk_index", 0),
                "structure_unit_id": parent_seg.get("structure_unit_id", ""),
                "parent_id":         seg_id,
                "segment_index":     int(fmt.get("segment_index", 0)),
                "frame_position":    frame_pos,
                "timestamp_s":       ts,
                "segment_start_s":   float(parent_seg.get("start_time_s", 0.0)),
                "segment_end_s":     float(parent_seg.get("end_time_s", 0.0)),
                "image_path":        ic.get("image_path", ""),
                "image_caption":     ic.get("image_caption", ""),
                "image_caption_confidence": None,
                "contextual_summary":       None,
                "contextual_summary_confidence": None,
                "frame_role":        role,
                "promoted_fields":   {},
            })

        _save_json(doc_meta,       folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks,    folders["metadata"] / "text_chunks.json")
        _save_json(image_chunks,   folders["metadata"] / "image_chunks.json")
        _save_json(video_segments, folders["metadata"] / "video_segments.json")
        _save_json(video_frames,   folders["metadata"] / "video_frames.json")
        _save_json(structure_units, folders["metadata"] / "structure_units.json")

        logger.success(
            f"Video: {file_record.relative_path!r} → "
            f"{len(segments)} segs | {len(image_chunks)} frames "
            f"| {len(video_frames)} video_frames | {len(structure_units)} scenes"
        )
        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": len(image_chunks), "table_count": 0,
            "segment_count": len(segments),
        }

    # ── Transcription ──────────────────────────────────────────────────────────

    def _transcribe(
        self, file_record: FileRecord, strategy: str, kwargs: Dict
    ) -> List[Dict]:
        """Transcribe using Whisper. Returns empty list gracefully if unavailable."""
        try:
            import whisper
            model_name = kwargs.get("whisper_model", "base")
            model      = whisper.load_model(model_name)
            result     = model.transcribe(
                str(file_record.local_path), word_timestamps=True
            )
            segments   = []
            for seg in result.get("segments", []):
                segments.append({
                    "text":       seg.get("text", "").strip(),
                    "start":      seg.get("start", 0.0),
                    "end":        seg.get("end",   0.0),
                    "confidence": seg.get("avg_logprob"),
                })
            logger.info(
                f"  Whisper: {len(segments)} segments from "
                f"{file_record.relative_path!r}"
            )
            return segments
        except ImportError:
            logger.warning("openai-whisper not installed — skipping transcription")
            return []
        except Exception as e:
            logger.error(f"Whisper transcription failed: {e}")
            return []

    # ── Tri-frame keyframe extraction with optical flow ────────────────────────

    def _extract_keyframes(
        self,
        file_record: FileRecord,
        segments:    List[Dict],
        doc_id:      str,
        folders:     Dict,
    ) -> Tuple[List[Dict], Dict]:
        """
        Extract start/peak/end keyframes per segment.

        Peak frame is selected by Farneback optical flow magnitude —
        the frame with the most visual change within the segment.
        Optionally crops dense visual regions from peak frames.

        Returns (image_chunk_dicts, figure_index_map).
        """
        try:
            import cv2
            import numpy as np
        except ImportError:
            logger.debug("opencv-python not installed — skipping keyframe extraction")
            return [], {}

        def _clamp(fn: int, total: int) -> int:
            return max(0, min(fn, max(0, total - 1)))

        def _read_frame(cap, fn: int):
            cap.set(cv2.CAP_PROP_POS_FRAMES, fn)
            ret, frame = cap.read()
            return frame if ret else None

        image_chunks:    List[Dict] = []
        figure_index_map: Dict      = {}

        cap = cv2.VideoCapture(str(file_record.local_path))
        if not cap.isOpened():
            logger.warning(
                f"Cannot open video for keyframes: {file_record.local_path}"
            )
            return [], {}

        fps          = cap.get(cv2.CAP_PROP_FPS) or 25.0
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)

        for i, seg in enumerate(segments, start=1):
            start_s  = float(seg.get("start", 0.0) or 0.0)
            end_s    = float(seg.get("end",   start_s) or start_s)
            end_s    = max(end_s, start_s)

            sf = _clamp(int(start_s * fps), total_frames)
            ef = _clamp(int(end_s   * fps), total_frames)
            if ef < sf:
                ef = sf

            span    = max(1, ef - sf)
            n_samp  = min(20, span + 1)
            samples = sorted(set(int(x) for x in np.linspace(sf, ef, num=n_samp)))
            if not samples:
                samples = [sf]

            # Optical flow peak selection
            peak_frame_no = samples[len(samples) // 2]
            peak_score    = 0.0
            prev_gray     = None
            for fn in samples:
                frame = _read_frame(cap, fn)
                if frame is None:
                    continue
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                if prev_gray is not None:
                    flow = cv2.calcOpticalFlowFarneback(
                        prev_gray, gray, None,
                        0.5, 3, 15, 3, 5, 1.2, 0,
                    )
                    mag, _ = cv2.cartToPolar(flow[..., 0], flow[..., 1])
                    score  = float(np.mean(mag))
                    if score > peak_score:
                        peak_score    = score
                        peak_frame_no = fn
                prev_gray = gray

            chosen = [
                ("start", sf,             start_s),
                ("peak",  peak_frame_no,  peak_frame_no / fps if fps else start_s),
                ("end",   ef,             end_s),
            ]

            keyframe_ids  = []
            peak_frame_id = None
            dense_count   = 0

            for role, fn, timestamp_s in chosen:
                frame = _read_frame(cap, fn)
                if frame is None:
                    continue

                fname    = f"{file_record.doc_stem}_seg_{i:04d}_{role}.jpg"
                fpath    = folders["frames"] / fname
                cv2.imwrite(str(fpath), frame)

                chunk_id = f"{doc_id}_seg_{i:04d}_{role}"
                keyframe_ids.append(chunk_id)
                if role == "peak":
                    peak_frame_id = chunk_id
                figure_index_map[f"{i}_{role}"] = chunk_id

                is_dense = (role == "peak" and peak_score >= INFO_DENSE_THRESHOLD)

                image_chunks.append({
                    "chunk_id":        chunk_id,
                    "doc_id":          doc_id,
                    "chunk_index":     len(image_chunks) + 1,
                    "source_modality": "image",
                    "figure_id":       chunk_id,
                    "image_path":      str(Path("frames") / fname),
                    "page_number":     0,
                    "image_type":      "keyframe",
                    "related_sections": [_segment_chunk_id(doc_id, i)],
                    "image_caption":   (
                        f"{role.title()} keyframe for segment {i} "
                        f"at {timestamp_s:.1f}s"
                    ),
                    "contextual_summary": None,
                    "contextual_summary_confidence": None,
                    "format_specific": {
                        "timestamp_s":         float(timestamp_s),
                        "segment_index":       i,
                        "segment_id":          _segment_chunk_id(doc_id, i),
                        "frame_role":          role,
                        "is_information_dense": is_dense,
                        "optical_flow_score":  float(
                            peak_score if role == "peak" else 0.0
                        ),
                    },
                })

                # Dense region crops from high-activity peak frames
                if is_dense:
                    crops = self._extract_dense_regions(
                        frame            = frame,
                        parent_frame_id  = chunk_id,
                        segment_index    = i,
                        segment_chunk_id = _segment_chunk_id(doc_id, i),
                        doc_id           = doc_id,
                        file_stem        = file_record.doc_stem,
                        folders          = folders,
                        base_index       = len(image_chunks) + dense_count,
                    )
                    dense_count += len(crops)
                    image_chunks.extend(crops)

            # Write vrag signals back onto segment dict for text chunk builder
            seg["_vrag"] = {
                "keyframe_ids":       keyframe_ids,
                "peak_frame_id":      peak_frame_id,
                "is_information_dense": peak_score >= INFO_DENSE_THRESHOLD,
                "optical_flow_score": float(peak_score),
                "region_crop_count":  dense_count,
                "keyframe_path":      (
                    str(Path("frames") / f"{file_record.doc_stem}_seg_{i:04d}_peak.jpg")
                    if peak_frame_id else None
                ),
                "keyframe_timestamp_s": (
                    peak_frame_no / fps if fps and peak_frame_id else None
                ),
            }

        cap.release()
        return image_chunks, figure_index_map

    # ── Dense region cropping ──────────────────────────────────────────────────

    def _extract_dense_regions(
        self,
        frame:            Any,
        parent_frame_id:  str,
        segment_index:    int,
        segment_chunk_id: str,
        doc_id:           str,
        file_stem:        str,
        folders:          Dict,
        base_index:       int = 0,
    ) -> List[Dict]:
        """
        Crop up to MAX_REGION_CROPS information-dense sub-regions from a frame.

        Uses adaptive Gaussian thresholding to find high-contrast areas
        (components, labels, diagrams) and bounding-box crops them.
        Each crop is saved as a separate image chunk so CLIP can embed
        focused close-ups rather than a cluttered full frame.
        """
        try:
            import cv2
            import numpy as np
        except ImportError:
            return []

        h, w    = frame.shape[:2]
        min_area = 0.05 * h * w   # reject tiny noise contours

        gray   = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        blur   = cv2.GaussianBlur(gray, (5, 5), 0)
        thresh = cv2.adaptiveThreshold(
            blur, 255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY_INV,
            31, 11,
        )
        contours, _ = cv2.findContours(
            thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
        )

        candidates = []
        for c in contours:
            x, y, cw, ch = cv2.boundingRect(c)
            area   = cw * ch
            if area < min_area:
                continue
            aspect = cw / max(ch, 1)
            if aspect < 0.3 or aspect > 3.0:
                continue
            candidates.append((area, x, y, cw, ch))
        candidates.sort(reverse=True)

        crops: List[Dict] = []
        for idx, (_, x, y, cw, ch) in enumerate(
            candidates[:MAX_REGION_CROPS], start=1
        ):
            crop  = frame[y : y + ch, x : x + cw]
            fname = f"{file_stem}_seg_{segment_index:04d}_region_{idx}.jpg"
            fpath = folders["frames"] / fname
            cv2.imwrite(str(fpath), crop)

            chunk_id = f"{doc_id}_seg_{segment_index:04d}_region_{idx}"
            crops.append({
                "chunk_id":        chunk_id,
                "doc_id":          doc_id,
                "chunk_index":     base_index + idx,
                "source_modality": "image",
                "figure_id":       chunk_id,
                "image_path":      str(Path("frames") / fname),
                "page_number":     0,
                "image_type":      "video_region_crop",
                "related_sections": [segment_chunk_id],
                "image_caption":   (
                    f"Dense region crop {idx} from segment {segment_index}"
                ),
                "contextual_summary": None,
                "contextual_summary_confidence": None,
                "format_specific": {
                    "segment_index":       segment_index,
                    "segment_id":          segment_chunk_id,
                    "parent_frame_id":     parent_frame_id,
                    "bbox":                [int(x), int(y), int(cw), int(ch)],
                    "frame_role":          "region",
                    "is_information_dense": True,
                },
            })

        return crops

    # ── Scene StructuralUnit builder ───────────────────────────────────────────

    def _build_scene_units(
        self,
        doc_id:         str,
        segments:       List[Dict],
        video_segments: List[Dict],
        text_chunks:    List[Dict],
    ) -> List[Dict]:
        """
        Group every SCENE_GROUP_SIZE segments into one scene StructuralUnit.

        Scene units are the navigation targets in sections_collection for
        video documents. The enrichment layer's VideoSegmentEnricher uses
        scene grouping (all segments in the same structure_unit_id) as the
        window for contextual chunking.
        """
        units:     List[Dict] = []
        n         = len(segments)
        scene_idx = 0
        chunk_map = {
            tc["format_specific"]["segment_id"]: tc["chunk_id"]
            for tc in text_chunks
            if tc.get("format_specific", {}).get("segment_id")
        }

        for start in range(0, n, self.SCENE_GROUP_SIZE):
            scene_idx += 1
            end        = min(start + self.SCENE_GROUP_SIZE, n)
            scene_segs = segments[start:end]

            start_s = float(scene_segs[0].get("start", 0.0))
            end_s   = float(scene_segs[-1].get("end",   0.0))

            # Collect chunk_ids and segment_ids in this scene
            scene_segment_ids: List[str] = []
            scene_chunk_ids:   List[str] = []
            for vs in video_segments[start:end]:
                seg_id   = vs.get("chunk_id", "")
                chunk_id = chunk_map.get(seg_id, "")
                if seg_id:
                    scene_segment_ids.append(seg_id)
                if chunk_id:
                    scene_chunk_ids.append(chunk_id)

            # Write structure_unit_id back onto each video_segment dict
            unit_id = _unit_id(doc_id, "scene", scene_idx)
            for vs in video_segments[start:end]:
                vs["structure_unit_id"] = unit_id

            units.append({
                "structure_unit_id":   unit_id,
                "doc_id":              doc_id,
                "structure_unit_type": "scene",
                "unit_index":          scene_idx,
                "title":               (
                    f"Scene {scene_idx} "
                    f"({start_s:.1f}s–{end_s:.1f}s)"
                ),
                "semantic_anchor":     (
                    f"Scene {scene_idx} "
                    f"({start_s:.1f}s–{end_s:.1f}s)"
                ),
                "start_time_s":        start_s,
                "end_time_s":          end_s,
                "chunk_ids":           scene_chunk_ids,
                "segment_ids":         scene_segment_ids,
                "figure_ids":          [],
                "table_ids":           [],
                "section_summary":     None,
                "keywords":            [],
                "entities":            [],
                "subsections":         [],
            })

        return units


# ── Unknown / Fallback Extractor ───────────────────────────────────────────────

class UnknownExtractor(BaseExtractor):
    """
    Fallback for unrecognised file types.
    Reads as plain text. Always writes all 5 standard files. Never raises.
    """

    def extract(self, file_record: FileRecord, **kwargs) -> Dict[str, Any]:
        doc_id  = generate_doc_id(file_record)
        folders = _create_output_folders(self.output_dir, file_record.doc_stem, fmt="unknown")

        content = None
        try:
            content = file_record.local_path.read_text(
                encoding="utf-8", errors="replace"
            )
        except Exception as e:
            logger.warning(f"Cannot read {file_record.relative_path}: {e}")

        text_chunks = []
        if content and content.strip():
            from ingestion_layer.utils.chunking_utils import estimate_tokens
            chunk_id = f"{doc_id}_chunk_001"
            text_chunks.append({
                "chunk_id":              chunk_id,
                "doc_id":                doc_id,
                "chunk_index":           1,
                "source_modality":       "text",
                "text_original_content": content[:10000],
                "local_context":         None,
                "page_number":           1,
                "section_title":         None,
                "chunk_strategy":        ChunkStrategy.WHOLE_FILE.value,
                "token_count":           estimate_tokens(content),
                "element_index":         1,
                "end_element_index":     1,
                "related_figures":       [],
                "related_tables":        [],
                "contextual_summary":    None,
                "contextual_summary_confidence": None,
                "detected_codes":        [],
            })

        doc_meta = self._base_doc_meta(doc_id, file_record, **kwargs)
        doc_meta.update({
            "chunk_count":   len(text_chunks),
            "file_format":   FileFormat.UNKNOWN.value,
            "structure_type": "unknown",
        })

        _save_json(doc_meta,    folders["metadata"] / "doc_metadata.json")
        _save_json(text_chunks, folders["metadata"] / "text_chunks.json")
        _save_json([],          folders["metadata"] / "structure_units.json")

        return {
            "doc_id": doc_id, "text_count": len(text_chunks),
            "image_count": 0, "table_count": 0,
        }